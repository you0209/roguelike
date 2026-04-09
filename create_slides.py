#!/usr/bin/env python3
"""DUNGEON QUEST ハッカソン発表スライド生成 (python-pptx版)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, subprocess

# ── Dimensions ────────────────────────────────────────────────────────────────
W = Inches(13.333)   # slide width
H = Inches(7.5)      # slide height
MX = Inches(0.5)     # side margin
CW = W - 2 * MX     # content width ≈ 12.333"
HDR_H = Inches(0.9)  # header bar height

# ── Colors ─────────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

BG   = rgb(15,  15,  35)   # #0f0f23 very dark navy
SURF = rgb(24,  30,  72)   # #181e48 dark blue surface
GOLD = rgb(255, 200, 50)   # #ffc832 gold
WHT  = rgb(255, 255, 255)
LITE = rgb(195, 210, 245)  # #c3d2f5 light blue-white
ACCB = rgb(100, 160, 255)  # #64a0ff blue accent
ACCG = rgb(90,  210, 120)  # #5ad278 green accent
ACCR = rgb(255, 110,  90)  # #ff6e5a red accent

# ── Helpers ────────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fc=None, lc=None, lw=1.5):
    """Add a colored rectangle shape."""
    s = slide.shapes.add_shape(1, x, y, w, h)   # 1 = RECTANGLE
    if fc:
        s.fill.solid()
        s.fill.fore_color.rgb = fc
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        fs=18, c=None, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font="Noto Sans JP",
        lspc=1.15, sa=0, sb=0):
    """Add a text box with styled text (newlines → multiple paragraphs)."""
    c = c or LITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Pt(0)
    tf.margin_right  = Pt(0)
    tf.margin_top    = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment    = align
        p.line_spacing = lspc
        if sa: p.space_before = Pt(sa)
        if sb: p.space_after  = Pt(sb)
        r = p.add_run()
        r.text       = line
        r.font.bold  = bold
        r.font.italic = italic
        r.font.size  = Pt(fs)
        r.font.color.rgb = c
        r.font.name  = font
    return tb

def hdr(slide, title):
    """Render the common gold header bar + slide title."""
    rect(slide, 0, 0, W, HDR_H, fc=GOLD)
    txt(slide, title,
        MX, Inches(0.1), CW, HDR_H - Inches(0.2),
        fs=25, c=BG, bold=True, align=PP_ALIGN.LEFT, lspc=1.0)

def card(slide, x, y, w, h, title, lines,
         tc=None, bc=None, tfs=17, bfs=14, lc=None):
    """Dark card: rectangle + title + body bullet lines."""
    tc = tc or GOLD
    bc = bc or LITE
    lc = lc or GOLD
    rect(slide, x, y, w, h, fc=SURF, lc=lc, lw=1.3)
    txt(slide, title,
        x + Inches(0.13), y + Inches(0.1),
        w - Inches(0.26), Inches(0.4),
        fs=tfs, c=tc, bold=True, align=PP_ALIGN.LEFT, lspc=1.0)
    txt(slide, "\n".join(lines),
        x + Inches(0.13), y + Inches(0.55),
        w - Inches(0.26), h - Inches(0.63),
        fs=bfs, c=bc, align=PP_ALIGN.LEFT, lspc=1.28, sa=2)

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

# ── Slide builders ─────────────────────────────────────────────────────────────

def s1_title(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    # Top / bottom gold bars
    rect(sl, 0, 0,            W, Inches(0.1), fc=GOLD)
    rect(sl, 0, H-Inches(0.1),W, Inches(0.1), fc=GOLD)
    # Center surface card
    rect(sl, MX, Inches(0.55), CW, Inches(4.0), fc=SURF)
    # Main title
    txt(sl, "DUNGEON QUEST",
        MX, Inches(0.75), CW, Inches(1.4),
        fs=72, c=GOLD, bold=True, align=PP_ALIGN.CENTER, font="Impact", lspc=1.0)
    # Subtitle
    txt(sl, "ローグライク × JRPG  ブラウザゲーム",
        MX, Inches(2.28), CW, Inches(0.62),
        fs=26, c=LITE, align=PP_ALIGN.CENTER, lspc=1.0)
    # Divider
    rect(sl, W // 3, Inches(3.06), W // 3, Inches(0.03), fc=GOLD)
    # Team members
    txt(sl, "慶田 優　／　齊藤 銀二　／　桑原 拓也",
        MX, Inches(3.18), CW, Inches(0.55),
        fs=22, c=WHT, align=PP_ALIGN.CENTER, lspc=1.0)
    # Workshop name
    txt(sl, "システム開発における生成AI活用ワークショップ",
        MX, H - Inches(0.88), CW, Inches(0.44),
        fs=16, c=ACCB, align=PP_ALIGN.CENTER, lspc=1.0)
    # Period
    txt(sl, "開発期間: 2026年 3月20日 〜 4月10日（約3週間）",
        MX, H - Inches(0.5),  CW, Inches(0.35),
        fs=14, c=ACCG, align=PP_ALIGN.CENTER, lspc=1.0)


def s2_agenda(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "目次  ―  Agenda")
    items = [
        ("01", "プロダクト概要（DUNGEON QUEST とは？）"),
        ("02", "なぜこのプロダクトを作ったか"),
        ("03", "生成 AI をどのように活用したか"),
        ("04", "実際の画面・デモ"),
        ("05", "こだわった点・工夫した点"),
        ("06", "開発で苦労したこと"),
        ("07", "今後の展望・改善点"),
    ]
    ih  = Inches(0.67)   # item height
    gap = Inches(0.07)
    by  = HDR_H + Inches(0.15)
    for i, (num, label) in enumerate(items):
        y = by + i * (ih + gap)
        rect(sl, MX, y, Inches(0.62), ih, fc=GOLD)
        txt(sl, num,
            MX + Inches(0.03), y + Inches(0.1),
            Inches(0.56), ih - Inches(0.2),
            fs=20, c=BG, bold=True, align=PP_ALIGN.CENTER, font="Impact", lspc=1.0)
        txt(sl, label,
            MX + Inches(0.72), y + Inches(0.12),
            CW - Inches(0.76), ih - Inches(0.22),
            fs=20, c=LITE, align=PP_ALIGN.LEFT, lspc=1.0)


def s3_overview(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "プロダクト概要 ―― DUNGEON QUEST とは？")
    # Description
    txt(sl,
        "7 階層のダンジョンを攻略し、最深部の龍神ヴァルムドラグを討伐するブラウザゲーム\n"
        "HTML / CSS / JavaScript のみで動作（外部ライブラリなし・index.html をブラウザで開くだけ）",
        MX, HDR_H + Inches(0.1), CW, Inches(0.72),
        fs=18, c=WHT, align=PP_ALIGN.CENTER, lspc=1.3)
    # 4 feature cards  (2 × 2)
    cw2 = (CW - Inches(0.18)) / 2
    ch  = Inches(2.55)
    by  = HDR_H + Inches(0.92)
    for i, (title, lines) in enumerate([
        ("   コマンドバトル",
         ["攻撃・防御・スキル・アイテム",
          "ターン制・プレイヤー先攻",
          "状態異常（毒・バフ / デバフ）",
          "ボスは専用スキルを使用"]),
        ("   レリックシステム",
         ["23 種のレリックを収集",
          "パッシブ効果が常時発動",
          "複数所持可・上限なし",
          "レア度で出現確率を調整"]),
        ("   豊富なイベント",
         ["各階層 5〜6 種 × 6 階層",
          "30 種超のランダムイベント",
          "選択肢でリスク vs リターン",
          "階層限定の特殊イベントも"]),
        ("   ショップ",
         ["階層移動のたびに入店",
          "レリック＋アイテムをランダム抽選",
          "ゴールドで強化・準備を整える",
          "戦略的な資源管理が重要"]),
    ]):
        cx = MX + (i % 2) * (cw2 + Inches(0.18))
        cy = by  + (i // 2) * (ch  + Inches(0.1))
        card(sl, cx, cy, cw2, ch, title, lines)


def s4_why(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "なぜこのプロダクトを作ったか")
    hw = (CW - Inches(0.18)) / 2
    card(sl, MX, HDR_H + Inches(0.15), hw, Inches(1.82),
         "   きっかけ・背景",
         ["ゲーム開発への興味と技術的挑戦",
          "ローグライクの高いリプレイ性 ×",
          "  JRPG のドラマ性を融合させたかった",
          "生成 AI でどこまで作れるか試したかった"],
         tfs=17, bfs=15)
    card(sl, MX + hw + Inches(0.18), HDR_H + Inches(0.15), hw, Inches(1.82),
         "   ワークショップ背景",
         ["千葉工大生向け特別企画として実施",
          "急速な生成 AI 技術の進化を受け",
          "「先端技術の開発の楽しさを学生に」",
          "  という思いから企画・実現"],
         tfs=17, bfs=15)
    card(sl, MX, HDR_H + Inches(2.12), CW, Inches(1.7),
         "   目指したゲーム体験・プロジェクト目標",
         ["外部ライブラリゼロで本格的なゲームを作りきる ―― 「作れる」という自信をつける",
          "限られたリソース管理 × ランダム性で毎回違う体験を実現する",
          "チーム 3 人の役割分担でリアルなチーム開発の流れを体験する",
          "「面白さ」を徹底的に考え抜き、バランス調整・アイデア出しに注力する"],
         tfs=18, bfs=15)
    txt(sl,
        "生成 AI を活用して、アイデア出しからプロトタイピング・実装・デモ発表まで\n"
        "一連の開発サイクルを約 3 週間で完走する実践的経験を積む",
        MX, HDR_H + Inches(3.97), CW, Inches(0.82),
        fs=17, c=ACCB, align=PP_ALIGN.CENTER, lspc=1.35)


def s5_ai(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "生成 AI をどのように活用したか")
    txt(sl, "Claude Code を全面採用 ―― コード・デバッグ・ドキュメントすべてを AI と協働で開発",
        MX, HDR_H + Inches(0.1), CW, Inches(0.5),
        fs=19, c=GOLD, bold=True, align=PP_ALIGN.CENTER, lspc=1.0)
    cw3 = (CW - Inches(0.24)) / 3
    ch  = Inches(4.55)
    by  = HDR_H + Inches(0.7)
    for i, (title, tc, lines) in enumerate([
        ("   コード生成", ACCG,
         ["全コードを Claude Code で生成",
          "バトルロジック・Canvas 描画",
          "ショップ・イベントシステム",
          "複雑なゲーム状態管理 (GS)",
          "19 種の敵スプライトも実装",
          "約 3000 行のコードを AI と協働で"]),
        ("   デバッグ支援", ACCB,
         ["バグ発生時に即 AI に相談",
          "UI の見切れ・コマンド非表示バグ",
          "バフのカウントダウンバグ",
          "レリック効果が機能しないバグ",
          "原因特定 → 修正まで一貫対応",
          "デバッグ時間を大幅に短縮"]),
        ("   ドキュメント", GOLD,
         ["README を AI が自動生成",
          "CHANGELOG を随時更新",
          "チーム分担表 HTML も生成",
          "設計方針・コメントも整備",
          "ドキュメント作成コストほぼゼロ",
          "チーム共有資料も AI が担当"]),
    ]):
        cx = MX + i * (cw3 + Inches(0.12))
        card(sl, cx, by, cw3, ch, title, lines, tc=tc, bfs=14, lc=tc)
    txt(sl, "AI なしでは難しかった ――「何が面白いか」のアイデア出し ／ ゲームバランスの方針決め",
        MX, by + ch + Inches(0.1), CW, Inches(0.42),
        fs=17, c=ACCG, align=PP_ALIGN.CENTER, lspc=1.0)


def s6_demo(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "実際の画面・デモ")
    # Demo placeholder area
    rect(sl, MX, HDR_H + Inches(0.12), CW, Inches(4.22), fc=SURF, lc=GOLD, lw=2.0)
    txt(sl, "【  デモ動画 / ライブデモ  】",
        MX + Inches(0.3), HDR_H + Inches(1.5), CW - Inches(0.6), Inches(1.5),
        fs=34, c=GOLD, bold=True, align=PP_ALIGN.CENTER, lspc=1.3)
    txt(sl, "index.html をブラウザで開くだけで動作します",
        MX + Inches(0.3), HDR_H + Inches(3.2), CW - Inches(0.6), Inches(0.6),
        fs=18, c=LITE, align=PP_ALIGN.CENTER, lspc=1.0)
    # Feature lists
    half = (CW - Inches(0.18)) / 2
    by   = HDR_H + Inches(4.45)
    txt(sl,
        "▶  タイトル → 階層選択（戦闘 / イベント / 階段）\n"
        "▶  コマンドバトル（スキル・アイテム使用）\n"
        "▶  ショップ・レリック購入\n"
        "▶  ランダムイベントの選択分岐",
        MX, by, half, Inches(1.75),
        fs=15, c=LITE, align=PP_ALIGN.LEFT, lspc=1.35, sa=2)
    txt(sl,
        "▶  龍神ヴァルムドラグとのボス戦\n"
        "▶  クリア・ゲームオーバー画面\n"
        "▶  スコア計算・表示\n"
        "▶  Canvas ドット絵スプライト 20 種",
        MX + half + Inches(0.18), by, half, Inches(1.75),
        fs=15, c=LITE, align=PP_ALIGN.LEFT, lspc=1.35, sa=2)


def s7_craft(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "こだわった点・工夫した点")
    cw2 = (CW - Inches(0.18)) / 2
    ch  = Inches(2.8)
    by  = HDR_H + Inches(0.15)
    for i, (title, lines) in enumerate([
        ("   ゲームバランスの徹底調整",
         ["同階層で戦闘を重ねるほど敵が強化（HP+12% / ATK+6%）",
          "7 階層の難易度上昇カーブを丁寧に設計",
          "レリック・アイテム価格のコストパフォーマンス調整"]),
        ("   30 種超のランダムイベント",
         ["各階層 5〜6 種 × 6 階層 = 30 種超",
          "選択肢でリスク vs リターンの面白さを実現",
          "階層限定イベントで後半の緊張感を演出"]),
        ("   23 種のレリックシステム",
         ["パッシブ効果が常時発動するビルド要素",
          "レア度（tier）で出現確率を重み付け",
          "複数所持でビルドの深みと戦略性を追加"]),
        ("   Canvas ドット絵スプライト",
         ["全 20 種のスプライトを Canvas 2D で手描き",
          "requestAnimationFrame によるアニメーション",
          "ボス専用スキルのエフェクト演出も実装"]),
    ]):
        cx = MX + (i % 2) * (cw2 + Inches(0.18))
        cy = by  + (i // 2) * (ch  + Inches(0.12))
        card(sl, cx, cy, cw2, ch, title, lines, tfs=18, bfs=15)


def s8_struggle(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "開発で苦労したこと")
    cw3 = (CW - Inches(0.24)) / 3
    ch  = Inches(5.5)
    by  = HDR_H + Inches(0.2)
    for i, (title, tc, lines) in enumerate([
        ("   ゲームバランスの調整", ACCR,
         ["「強すぎず・弱すぎず」の絶妙な",
          "数値設計が想像以上に難しかった",
          "",
          "数値を少し変えるだけで体験が",
          "大きく変化し、何度も試行錯誤",
          "",
          "AI に相談しながら",
          "バランスの方針を決定した",
          "",
          "プレイテストを繰り返して",
          "少しずつ理想に近づけた"]),
        ("   アイデア出し", ACCB,
         ["「何が面白いか」の答えは",
          "数値では出ない創造的な作業",
          "",
          "イベントやレリックの効果案を",
          "考えるのが最も難しかった",
          "",
          "AI へのプロンプトを工夫して",
          "アイデアの壁を突破した",
          "",
          "「面白さ」を言語化する",
          "訓練になった"]),
        ("   チーム開発のマージ", ACCG,
         ["3 人が同時に異なる機能を実装",
          "コンフリクトを防ぐため",
          "ファイルを機能ごとに分割",
          "",
          "担当ファイルの境界を明確に",
          "共通ファイル変更時は事前連絡",
          "ブランチ戦略で安全に統合",
          "",
          "チーム開発の難しさと",
          "醍醐味を同時に体験した"]),
    ]):
        cx = MX + i * (cw3 + Inches(0.12))
        rect(sl, cx, by, cw3, ch, fc=SURF, lc=tc, lw=2.0)
        txt(sl, title,
            cx + Inches(0.12), by + Inches(0.1),
            cw3 - Inches(0.24), Inches(0.4),
            fs=16, c=tc, bold=True, align=PP_ALIGN.LEFT, lspc=1.0)
        txt(sl, "\n".join(lines),
            cx + Inches(0.12), by + Inches(0.56),
            cw3 - Inches(0.24), ch - Inches(0.65),
            fs=14, c=LITE, align=PP_ALIGN.LEFT, lspc=1.32, sa=2)


def s9_future(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "今後の展望・改善点")
    hw = (CW - Inches(0.18)) / 2
    ch = Inches(5.5)
    by = HDR_H + Inches(0.2)
    card(sl, MX, by, hw, ch,
         "   近い将来の改善",
         ["宝箱イベントにアイテム報酬追加",
          "  （4〜6 階の宝箱を充実させる）",
          "",
          "スキル入手方法の見直し",
          "  （ショップタブ廃止後の代替手段）",
          "",
          "BGM / SE の追加",
          "  （サウンドなしは没入感が薄い）",
          "",
          "仲間システムの実装",
          "  （パーティー編成要素の追加）"],
         tc=ACCG, tfs=18, bfs=15, lc=ACCG)
    card(sl, MX + hw + Inches(0.18), by, hw, ch,
         "   将来の展望",
         ["スマートフォン対応（タッチ UI）",
          "",
          "オンラインスコアランキング",
          "",
          "より多くのボス・フロアの追加",
          "",
          "AI 生成によるダンジョンマップの",
          "  動的生成への挑戦",
          "",
          "「生成 AI × ゲーム生成」",
          "  プロジェクトへの発展"],
         tc=ACCB, tfs=18, bfs=15, lc=ACCB)


def s10_summary(prs):
    sl = new_slide(prs)
    set_bg(sl, BG)
    hdr(sl, "まとめ")
    rect(sl, 0, H - Inches(0.1), W, Inches(0.1), fc=GOLD)
    txt(sl, "DUNGEON QUEST  ――  ローグライク × JRPG  ブラウザゲーム",
        MX, HDR_H + Inches(0.1), CW, Inches(0.55),
        fs=23, c=GOLD, bold=True, align=PP_ALIGN.CENTER, lspc=1.0)
    for i, (text, c) in enumerate([
        ("Claude Code を全面活用し、3 人チームで 3 週間でゲームを完成させた",      ACCG),
        ("コード生成・デバッグ・ドキュメントすべてを AI と協働で実現できた",        ACCB),
        ("「面白いゲームを作る」という本質的な課題に集中することができた",          LITE),
        ("生成 AI が開発の可能性を大きく広げることを体験・実感した",              WHT),
    ]):
        y = HDR_H + Inches(0.82) + i * Inches(1.03)
        rect(sl, MX, y, Inches(0.55), Inches(0.5), fc=GOLD)
        txt(sl, "✓",
            MX + Inches(0.03), y + Inches(0.04), Inches(0.49), Inches(0.44),
            fs=22, c=BG, bold=True, align=PP_ALIGN.CENTER, lspc=1.0)
        txt(sl, text,
            MX + Inches(0.65), y + Inches(0.08), CW - Inches(0.68), Inches(0.42),
            fs=20, c=c, align=PP_ALIGN.LEFT, lspc=1.0)
    txt(sl, "ご清聴ありがとうございました",
        MX, H - Inches(1.1), CW, Inches(0.58),
        fs=28, c=GOLD, bold=True, align=PP_ALIGN.CENTER, lspc=1.0)
    txt(sl, "慶田 優　／　齊藤 銀二　／　桑原 拓也",
        MX, H - Inches(0.58), CW, Inches(0.36),
        fs=15, c=LITE, align=PP_ALIGN.CENTER, lspc=1.0)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    builders = [s1_title, s2_agenda, s3_overview, s4_why,
                s5_ai, s6_demo, s7_craft, s8_struggle,
                s9_future, s10_summary]
    for builder in builders:
        builder(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "dungeon_quest_slides.pptx")
    prs.save(out)
    print(f"✅ 完成: {out}")
    return out


if __name__ == '__main__':
    path = main()
    subprocess.run(["open", path])   # macOS: PowerPoint or Google Slides で開く
